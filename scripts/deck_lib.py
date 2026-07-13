#!/usr/bin/env python3
"""deck_lib.py — layout + token library for /vibe:deck (Ink & Signal system).

Import from a per-deck build script run with the vibe-pptx toolkit venv:
    import sys; sys.path.insert(0, "<plugin>/scripts")
    from deck_lib import Deck, TOKENS

Tokens (TOKENS / TYPE / SEM / GANTT_PHASES) are cited to DR rules in
references/design-system.md and mirrored as hard gates in validate_deck.py —
change a threshold in the research report first, then here and there in lockstep.

Four quality levers (from the deck-taste research + the devant study):
  1. Ink & Signal palette (warm near-black, one signal accent, derived neutrals)
  2. Semantic typed diagram nodes (action/decision/success/error/data/external)
  3. Taste-move archetypes (last-card-accent, accent kicker mark, chevrons,
     accent-rule three-point, hero metric that needs a source)
  4. anti-slop discipline (single-accent chrome, shadowless, one font, no em-dash)

Everything is drawn as native shapes on the standard 13.333x7.5in canvas, so
the build stays instant (no LibreOffice --convert-to in the build path).
"""
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# Canvas + grid (standard PowerPoint 16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Emu(12192000)   # 13.333in
SLIDE_H = Emu(6858000)    # 7.5in
MARGIN = Inches(0.9)
CONTENT_W = Inches(13.333 - 1.8)
CONTENT_R = Inches(13.333 - 0.9)


# ---------------------------------------------------------------------------
# Ink & Signal palette (devant brand.json values, exact)
# ---------------------------------------------------------------------------
def _hex(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lerp(a, b, t):
    return RGBColor(*(round(a[i] + (b[i] - a[i]) * t) for i in range(3)))


# derive() reproduces the neutral ramp from just ink/paper/accent (the
# "rebrand from 3 hexes" story); the demo uses devant's exact published hexes.
def derive(ink, paper, accent):
    return {
        "ink": ink, "paper": paper, "accent": accent,
        "panel": _lerp(paper, ink, 0.06),
        "hairline": _lerp(paper, ink, 0.12),
        "mist": _lerp(paper, ink, 0.40),
        "grey": _lerp(paper, ink, 0.58),
        "accentTint": _lerp(accent, paper, 0.78),
    }


TOKENS = {
    "ink": _hex("16283F"),        # deep navy (dark backgrounds + text)
    "paper": _hex("FFFFFF"),
    "panel": _hex("EEF2F7"),      # cool light blue-grey band
    "grey": _hex("5C6675"),       # captions/secondary on light
    "mist": _hex("93A1B2"),       # secondary on navy
    "hairline": _hex("DCE3EC"),   # dividers / table rules
    "accent": _hex("1F5FA8"),     # THE one signal accent (blue)
    "accentTint": _hex("D7E3F1"),
}

# Semantic diagram palette (devant draw.io swatches, colour-blind-safe).
# Allowed ONLY inside diagram nodes — deck chrome stays single-accent.
SEM = {
    "action":   (_hex("DAE8FC"), _hex("6C8EBF")),
    "success":  (_hex("D5E8D4"), _hex("82B366")),
    "error":    (_hex("F8CECC"), _hex("B85450")),
    "decision": (_hex("FFF2CC"), _hex("D6B656")),
    "data":     (_hex("E1D5E7"), _hex("9673A6")),
    "external": (_hex("F5F5F5"), _hex("666666")),
}

# Roadmap/Gantt month ramp (light -> dark navy), matches the reference image.
GANTT_PHASES = [_hex("3E7CBF"), _hex("1F5FA8"), _hex("123A66")]

FONT = "Noto Sans"   # devant brand font (deliberately not Inter/Arial)

# Type scale (devant _PARA, adapted +1-2pt for the larger standard canvas).
TYPE = {
    "fig": Pt(80),      # hero figure, accent
    "h1": Pt(42),       # cover / statement headline
    "h2": Pt(28),       # slide / section headline
    "h2w": Pt(25),      # headline on dark
    "num": Pt(23),      # card numbers
    "chev": Pt(22),     # gutter chevrons
    "sub": Pt(18),      # subtitle
    "card": Pt(16),     # card titles
    "body": Pt(15),     # body / node labels
    "eye": Pt(14),      # eyebrow / kicker (UPPERCASE, tracked)
    "cap": Pt(14),      # captions / sources / node notes (DK-2 14pt floor)
}


# ---------------------------------------------------------------------------
# low-level helpers
# ---------------------------------------------------------------------------
def _flatten(shape):
    """Drop <p:style> so no theme effectRef (drop shadow) survives."""
    style = shape._element.find(qn("p:style"))
    if style is not None:
        shape._element.remove(style)


def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, fill, line=None, w=Pt(1.25)):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shape)
    else:
        shape.line.color.rgb = line
        shape.line.width = w
    _flatten(shape)


def _set_spc(run, pts):
    """Character spacing (letter-tracking) in points, via rPr@spc (1/100 pt)."""
    run.font._rPr.set("spc", str(int(pts * 100)))


def _dash(conn, kind="dash"):
    ln = conn.line._get_or_add_ln()
    d = ln.find(qn("a:prstDash"))
    if d is None:
        d = ln.makeelement(qn("a:prstDash"), {})
        ln.append(d)
    d.set("val", kind)


def _arrow(conn, end="tail", size="med"):
    """Add a triangular arrowhead. head=begin, tail=end."""
    ln = conn.line._get_or_add_ln()
    tag = "a:tailEnd" if end == "tail" else "a:headEnd"
    if ln.find(qn(tag)) is None:
        ln.append(ln.makeelement(qn(tag), {"type": "triangle", "w": size, "len": size}))


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]

    def save(self, path):
        self.prs.save(path)

    # -- slide + primitives ------------------------------------------------
    def _slide(self, bg=None):
        s = self.prs.slides.add_slide(self._blank)
        if bg is not None:
            r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
            _solid(r, bg)
            el = r._element
            el.getparent().remove(el)
            s.shapes._spTree.insert(2, el)  # send to back
        return s

    def _run(self, para, text, size, color, bold=False, italic=False, spc=None):
        r = para.add_run()
        r.text = text
        r.font.size = size
        r.font.name = FONT
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if spc is not None:
            _set_spc(r, spc)
        return r

    def _text(self, slide, x, y, w, h, name, lines, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
        """lines: list of dicts {t,size,color,bold,italic,spc,upper,align,space_after}."""
        box = slide.shapes.add_textbox(x, y, w, h)
        box.name = name
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = ln.get("align", align)
            if ln.get("space_after") is not None:
                p.space_after = ln["space_after"]
            txt = ln["t"].upper() if ln.get("upper") else ln["t"]
            self._run(p, txt, ln["size"], ln["color"], ln.get("bold", False),
                      ln.get("italic", False), ln.get("spc"))
        return box

    def _kicker(self, slide, text, y=Inches(0.72), dark=False, mark=True):
        """Accent square mark + tracked uppercase eyebrow."""
        if mark:
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y,
                                        Inches(0.18), Inches(0.18))
            _solid(sq, TOKENS["accent"])
            tx = Emu(int(MARGIN) + Inches(0.32))
        else:
            tx = MARGIN
        self._text(slide, tx, Emu(int(y) - int(Inches(0.03))), CONTENT_W,
                   Inches(0.4), "kicker",
                   [{"t": text, "size": TYPE["eye"], "color": TOKENS["accent"],
                     "bold": True, "upper": True, "spc": 2.0}])

    def _rule(self, slide, x, y, w, color=None, h=Inches(0.055)):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        _solid(r, color or TOKENS["accent"])
        return r

    def _footer(self, slide, text, dark=False):
        self._text(slide, MARGIN, Inches(7.02), CONTENT_W, Inches(0.35),
                   "footer",
                   [{"t": text, "size": TYPE["cap"],
                     "color": TOKENS["mist"] if dark else TOKENS["grey"]}])

    def _heading(self, slide, kicker, heading):
        y = Inches(0.72)
        if kicker:
            self._kicker(slide, kicker, y)
            hy = Inches(1.18)
        else:
            hy = Inches(0.85)
        self._text(slide, MARGIN, hy, CONTENT_W, Inches(1.1), "title",
                   [{"t": heading, "size": TYPE["h2"], "color": TOKENS["ink"],
                     "bold": True}])
        return slide

    # -- diagram node ------------------------------------------------------
    def _node(self, slide, x, y, w, h, sem, label, note=None, name="node",
              text_dark=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        fill, line = SEM[sem]
        b = slide.shapes.add_shape(shape, x, y, w, h)
        b.name = name
        _solid(b, fill, line, Pt(1.75))
        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, label, TYPE["body"], TOKENS["ink"], bold=True)
        if note:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            self._run(p2, note, TYPE["cap"], TOKENS["grey"])
        return b

    def _connect(self, slide, a, b, a_side, b_side, color=None, dash=None,
                 kind=MSO_CONNECTOR.STRAIGHT, w=Pt(1.75)):
        c = slide.shapes.add_connector(kind, 0, 0, 0, 0)
        c.begin_connect(a, a_side)
        c.end_connect(b, b_side)
        c.line.color.rgb = color or TOKENS["grey"]
        c.line.width = w
        _flatten(c)
        if dash:
            _dash(c, dash)
        _arrow(c)
        return c

    def _free_arrow(self, slide, pts, color=None, dash=None, w=Pt(1.75)):
        """Orthogonal poly-line arrow through explicit (x,y) Inches points."""
        from pptx.util import Emu as _E
        xs = [int(p[0]) for p in pts]
        ys = [int(p[1]) for p in pts]
        left, top = min(xs), min(ys)
        width = max(max(xs) - left, 1)
        height = max(max(ys) - top, 1)
        fb = slide.shapes.build_freeform(_E(xs[0] - left), _E(ys[0] - top), scale=1.0)
        fb.add_line_segments([(_E(x - left), _E(y - top)) for x, y in zip(xs[1:], ys[1:])],
                             close=False)
        shp = fb.convert_to_shape()
        shp.left, shp.top, shp.width, shp.height = _E(left), _E(top), _E(width), _E(height)
        shp.fill.background()
        shp.line.color.rgb = color or TOKENS["grey"]
        shp.line.width = w
        _flatten(shp)
        if dash:
            _dash(shp, dash)
        _arrow(shp)
        return shp

    def _legend(self, slide, items, y=Inches(6.75)):
        """items: list of (sem, label). Small swatch + label row, bottom-left."""
        x = MARGIN
        for sem, label in items:
            fill, line = SEM[sem]
            sw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                                        Inches(0.22), Inches(0.16))
            _solid(sw, fill, line, Pt(1.25))
            self._text(slide, Emu(int(x) + int(Inches(0.30))),
                       Emu(int(y) - int(Inches(0.05))), Inches(2.2), Inches(0.3),
                       "legend", [{"t": label, "size": TYPE["cap"],
                                   "color": TOKENS["grey"]}])
            x = Emu(int(x) + int(Inches(0.30)) + int(Inches(1.1 + 0.02 * len(label))))
        return slide

    # =====================================================================
    # ARCHETYPES
    # =====================================================================
    def title(self, headline, kicker=None, subtitle=None, footer=None, eyebrow=None):
        kicker = kicker or eyebrow  # back-compat alias
        s = self._slide(bg=TOKENS["ink"])
        sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.55),
                                Inches(0.34), Inches(0.34))
        _solid(sq, TOKENS["accent"])
        if kicker:
            self._text(s, Emu(int(MARGIN) + int(Inches(0.52))), Inches(2.55),
                       CONTENT_W, Inches(0.4), "kicker",
                       [{"t": kicker, "size": TYPE["eye"], "color": TOKENS["accent"],
                         "bold": True, "upper": True, "spc": 2.0}])
        self._text(s, MARGIN, Inches(3.15), CONTENT_W, Inches(2.0), "title",
                   [{"t": headline, "size": TYPE["h1"], "color": TOKENS["paper"],
                     "bold": True}])
        self._rule(s, MARGIN, Inches(5.2), Inches(3.3), TOKENS["accent"])
        if subtitle:
            self._text(s, MARGIN, Inches(5.45), CONTENT_W, Inches(0.9), "subtitle",
                       [{"t": subtitle, "size": TYPE["sub"], "color": TOKENS["mist"]}])
        # nested-loop motif (the deck's theme: loops within loops), top-right
        self._glyph(s, int(Inches(9.55)), int(Inches(0.95)), int(Inches(2.5)),
                    MSO_SHAPE.CIRCULAR_ARROW, TOKENS["mist"], w=Pt(2.25))
        self._glyph(s, int(Inches(10.25)), int(Inches(1.62)), int(Inches(1.15)),
                    MSO_SHAPE.CIRCULAR_ARROW, TOKENS["accent"], w=Pt(2.25))
        if footer:
            self._text(s, MARGIN, Inches(6.85), CONTENT_W, Inches(0.4), "footer",
                       [{"t": footer, "size": TYPE["cap"], "color": TOKENS["mist"]}])
        return s

    def section(self, number, title):
        s = self._slide(bg=TOKENS["ink"])
        self._rule(s, MARGIN, Inches(2.4), Inches(3.3), TOKENS["accent"])
        self._text(s, MARGIN, Inches(2.7), CONTENT_W, Inches(0.5), "kicker",
                   [{"t": number, "size": TYPE["eye"], "color": TOKENS["accent"],
                     "bold": True, "upper": True, "spc": 2.0}])
        self._text(s, MARGIN, Inches(3.25), CONTENT_W, Inches(1.6), "title",
                   [{"t": title, "size": TYPE["h1"], "color": TOKENS["paper"],
                     "bold": True}])
        return s

    def statement(self, text, source=None, kicker=None, eyebrow=None):
        kicker = kicker or eyebrow  # back-compat alias
        s = self._slide(bg=TOKENS["paper"])
        self._rule(s, MARGIN, Inches(2.1), Inches(0.9), TOKENS["accent"])
        if kicker:
            self._text(s, MARGIN, Inches(1.5), CONTENT_W, Inches(0.4), "kicker",
                       [{"t": kicker, "size": TYPE["eye"], "color": TOKENS["accent"],
                         "bold": True, "upper": True, "spc": 2.0}])
        self._text(s, MARGIN, Inches(2.5), CONTENT_W, Inches(2.9), "statement",
                   [{"t": text, "size": TYPE["h1"], "color": TOKENS["ink"],
                     "bold": True}], anchor=MSO_ANCHOR.TOP)
        if source:
            self._text(s, MARGIN, Inches(6.2), CONTENT_W, Inches(0.5), "source",
                       [{"t": source, "size": TYPE["cap"], "color": TOKENS["grey"],
                         "italic": True}])
        return s

    def cards(self, heading, items, kicker=None, last_accent=True, source=None,
              eyebrow=None):
        """items: list of (title, desc). Numbered; last card filled accent."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker or eyebrow, heading)
        n = len(items)
        gap = Inches(0.45)
        chev_w = Inches(0.35)
        card_w = Emu(int((int(CONTENT_W) - int(gap) * (n - 1)) / n))
        top = Inches(2.55)
        card_h = Inches(3.5)
        x = MARGIN
        for i, (title, desc) in enumerate(items):
            is_last = last_accent and i == n - 1
            fill = TOKENS["accent"] if is_last else TOKENS["panel"]
            tcol = TOKENS["paper"] if is_last else TOKENS["ink"]
            dcol = TOKENS["accentTint"] if is_last else TOKENS["grey"]
            ncol = TOKENS["paper"] if is_last else TOKENS["accent"]
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, card_h)
            card.name = f"card{i}"
            _solid(card, fill, None if is_last else TOKENS["hairline"], Pt(1.0))
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.28)
            tf.margin_top = Inches(0.28)
            p0 = tf.paragraphs[0]
            self._run(p0, f"{i + 1:02d}", TYPE["num"], ncol, bold=True)
            p1 = tf.add_paragraph()
            p1.space_before = Pt(10)
            self._run(p1, title, TYPE["card"], tcol, bold=True)
            p2 = tf.add_paragraph()
            p2.space_before = Pt(6)
            self._run(p2, desc, TYPE["body"], dcol)
            _flatten(card)
            if i < n - 1:
                cx = Emu(int(x) + int(card_w) + int((int(gap) - int(chev_w)) / 2))
                self._text(s, cx, Emu(int(top) + int(Inches(1.4))), chev_w,
                           Inches(0.6), "chev",
                           [{"t": ">", "size": TYPE["chev"], "color": TOKENS["mist"],
                             "bold": True, "align": PP_ALIGN.CENTER}])
            x = Emu(int(x) + int(card_w) + int(gap))
        if source:
            self._footer(s, source)
        return s

    def two_column(self, heading, left, right, kicker=None, source=None):
        """left/right: dict {title, points:[...]}. Right card is accent."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, heading)
        gap = Inches(0.6)
        col_w = Emu(int((int(CONTENT_W) - int(gap)) / 2))
        top = Inches(2.55)
        col_h = Inches(3.5)
        for i, (data, accent) in enumerate([(left, False), (right, True)]):
            x = MARGIN if i == 0 else Emu(int(MARGIN) + int(col_w) + int(gap))
            fill = TOKENS["accent"] if accent else TOKENS["panel"]
            tcol = TOKENS["paper"] if accent else TOKENS["ink"]
            bcol = TOKENS["accentTint"] if accent else TOKENS["grey"]
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
            card.name = f"col{i}"
            _solid(card, fill, None if accent else TOKENS["hairline"], Pt(1.0))
            tf = card.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.36)
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT
            self._run(p0, data["title"], TYPE["card"], tcol, bold=True)
            for pt in data["points"]:
                pp = tf.add_paragraph()
                pp.alignment = PP_ALIGN.LEFT
                pp.space_before = Pt(11)
                self._run(pp, pt, TYPE["body"], bcol)
            _flatten(card)
        if source:
            self._footer(s, source)
        return s

    def three_point(self, heading, points, kicker=None, source=None):
        """points: list of (title, desc). Accent rule per column, no card fill."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, heading)
        n = len(points)
        gap = Inches(0.6)
        col_w = Emu(int((int(CONTENT_W) - int(gap) * (n - 1)) / n))
        top = Inches(2.7)
        x = MARGIN
        for title, desc in points:
            self._rule(s, x, top, Inches(0.7), TOKENS["accent"])
            self._text(s, x, Emu(int(top) + int(Inches(0.28))), col_w, Inches(0.7),
                       "pt-title", [{"t": title, "size": TYPE["card"],
                                     "color": TOKENS["ink"], "bold": True}])
            self._text(s, x, Emu(int(top) + int(Inches(1.05))), col_w, Inches(2.6),
                       "pt-desc", [{"t": desc, "size": TYPE["body"],
                                    "color": TOKENS["grey"]}])
            x = Emu(int(x) + int(col_w) + int(gap))
        if source:
            self._footer(s, source)
        return s

    def metric(self, value, label, src, heading=None, kicker=None):
        s = self._slide(bg=TOKENS["ink"])
        if kicker:
            self._text(s, MARGIN, Inches(0.72), CONTENT_W, Inches(0.4), "kicker",
                       [{"t": kicker, "size": TYPE["eye"], "color": TOKENS["accent"],
                         "bold": True, "upper": True, "spc": 2.0}])
        self._text(s, MARGIN, Inches(2.1), CONTENT_W, Inches(2.0), "figure",
                   [{"t": value, "size": TYPE["fig"], "color": TOKENS["accent"],
                     "bold": True}])
        self._text(s, MARGIN, Inches(4.2), CONTENT_W, Inches(0.8), "label",
                   [{"t": label, "size": TYPE["sub"], "color": TOKENS["paper"]}])
        self._text(s, MARGIN, Inches(6.85), CONTENT_W, Inches(0.4), "source",
                   [{"t": src, "size": TYPE["cap"], "color": TOKENS["mist"],
                     "italic": True}])
        return s

    def big_number(self, value, context, title=None, source=None, kicker=None):
        """One hero metric (~6x body) on a light slide + one context line."""
        s = self._slide(bg=TOKENS["paper"])
        if title:
            self._heading(s, kicker, title)
        self._text(s, MARGIN, Inches(2.5), CONTENT_W, Inches(2.2), "hero",
                   [{"t": value, "size": TYPE["fig"], "color": TOKENS["accent"],
                     "bold": True}])
        self._text(s, MARGIN, Inches(4.9), CONTENT_W, Inches(1.0), "context",
                   [{"t": context, "size": TYPE["body"], "color": TOKENS["ink"]}])
        if source:
            self._footer(s, source)
        return s

    def title_body(self, title, body_lines, eyebrow=None, kicker=None, source=None):
        """Assertion title + up to ~4 short body lines (never a bullet wall)."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker or eyebrow, title)
        lines = [{"t": ln, "size": TYPE["body"], "color": TOKENS["ink"],
                  "space_after": Pt(9)} for ln in body_lines]
        self._text(s, MARGIN, Inches(2.6), CONTENT_W, Inches(3.9), "body", lines)
        if source:
            self._footer(s, source)
        return s

    def matrix(self, heading, columns, rows, emphasize=(), kicker=None,
               source=None, col_weights=None):
        """Status table from native rects.
        columns: list of header strings. rows: list of [cell,...] (len==columns).
        emphasize: set of row indices drawn with an accent left-edge + bold.
        """
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, heading)
        top = Inches(2.35)
        table_w = int(CONTENT_W)
        ncol = len(columns)
        weights = col_weights or [1.0] * ncol
        wsum = sum(weights)
        col_x = [MARGIN]
        for w in weights[:-1]:
            col_x.append(Emu(int(col_x[-1]) + int(table_w * w / wsum)))
        col_px = [int(table_w * w / wsum) for w in weights]
        n = len(rows)
        head_h = Inches(0.5)
        avail = int(Inches(6.7)) - int(top) - int(head_h)
        row_h = Emu(int(avail / max(n, 1)))
        # header band
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top, Emu(table_w), head_h)
        _solid(band, TOKENS["panel"])
        for j, col in enumerate(columns):
            self._text(s, Emu(int(col_x[j]) + int(Inches(0.12))),
                       Emu(int(top) + int(Inches(0.09))), Emu(col_px[j] - int(Inches(0.2))),
                       Inches(0.4), "cell-h",
                       [{"t": col, "size": TYPE["cap"], "color": TOKENS["grey"],
                         "bold": True, "upper": True, "spc": 0.8}])
        y = Emu(int(top) + int(head_h))
        for i, row in enumerate(rows):
            emph = i in emphasize
            if emph:
                hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, Emu(table_w), row_h)
                _solid(hl, TOKENS["paper"])
                edge = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y,
                                          Inches(0.06), row_h)
                _solid(edge, TOKENS["accent"])
            for j, cell in enumerate(row):
                col = TOKENS["ink"] if (emph or j == 0) else TOKENS["grey"]
                self._text(s, Emu(int(col_x[j]) + int(Inches(0.16))),
                           Emu(int(y) + int(Inches(0.12))),
                           Emu(col_px[j] - int(Inches(0.28))), row_h, "cell",
                           [{"t": cell, "size": TYPE["body"], "color": col,
                             "bold": emph and j == 0}])
            rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN,
                                      Emu(int(y) + int(row_h)), Emu(table_w),
                                      Inches(0.012))
            _solid(rule, TOKENS["hairline"])
            y = Emu(int(y) + int(row_h))
        if source:
            self._footer(s, source)
        return s

    def closing(self, text, contact=None):
        s = self._slide(bg=TOKENS["ink"])
        self._glyph(s, int(MARGIN), int(Inches(1.55)), int(Inches(1.0)),
                    MSO_SHAPE.CIRCULAR_ARROW, TOKENS["accent"], w=Pt(3.0))
        self._rule(s, MARGIN, Inches(2.6), Inches(0.9), TOKENS["accent"])
        self._text(s, MARGIN, Inches(3.0), CONTENT_W, Inches(2.0), "statement",
                   [{"t": text, "size": TYPE["h1"], "color": TOKENS["paper"],
                     "bold": True}])
        if contact:
            self._text(s, MARGIN, Inches(5.3), CONTENT_W, Inches(0.5), "contact",
                       [{"t": contact, "size": TYPE["sub"], "color": TOKENS["mist"]}])
        return s

    # =====================================================================
    # DIAGRAM ARCHETYPES (semantic native shapes)
    # =====================================================================
    def _polyline(self, slide, pts, color, dash=None, arrow=True, w=Pt(1.75)):
        """Orthogonal arrow through explicit EMU (x,y) points, as line segments."""
        segs = []
        for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
            c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                           Emu(x1), Emu(y1), Emu(x2), Emu(y2))
            c.line.color.rgb = color
            c.line.width = w
            _flatten(c)
            if dash:
                _dash(c, dash)
            segs.append(c)
        if arrow and segs:
            _arrow(segs[-1])
        return segs

    def _row(self, n, left, right, gap):
        """n equal cells between left/right (EMU), equal gaps → list of (x,w)."""
        span = right - left
        w = int((span - gap * (n - 1)) / n)
        return [(left + i * (w + gap), w) for i in range(n)]

    def flow(self, heading, nodes, kicker=None, backedge=None, source=None,
             legend=None, eyebrow=None, y=Inches(3.35), node_h=Inches(1.15)):
        """Semantic left-to-right flow. nodes: list of (sem, label[, note]);
        a bare string is treated as an ("action", label) node.
        backedge: dict {label, crossed(bool), color, from_idx, to_idx}."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker or eyebrow, heading)
        nodes = [("action", it) if isinstance(it, str) else it for it in nodes]
        n = len(nodes)
        gap = int(Inches(0.4))
        cells = self._row(n, int(MARGIN), int(CONTENT_R), gap)
        boxes = []
        for (x, w), item in zip(cells, nodes):
            sem, label = item[0], item[1]
            note = item[2] if len(item) > 2 else None
            shape = MSO_SHAPE.DIAMOND if sem == "decision" else MSO_SHAPE.ROUNDED_RECTANGLE
            boxes.append(self._node(s, Emu(x), y, Emu(w), node_h, sem, label,
                                    note, name="node", shape=shape))
        for a, b in zip(boxes, boxes[1:]):
            self._connect(s, a, b, 3, 1, color=TOKENS["grey"])
        if backedge:
            crossed = backedge.get("crossed")
            col = backedge.get("color") or (SEM["error"][1] if crossed else TOKENS["accent"])
            first = boxes[backedge.get("to_idx", 0)]
            last = boxes[backedge.get("from_idx", n - 1)]
            fx = int(first.left) + int(first.width) // 2
            lx = int(last.left) + int(last.width) // 2
            by = int(y) + int(node_h)
            lane = by + int(Inches(0.7))
            self._polyline(s, [(lx, by), (lx, lane), (fx, lane), (fx, by)],
                           col, dash="dash")
            mid = (fx + lx) // 2
            if crossed:
                d = int(Inches(0.4))
                sym = s.shapes.add_shape(MSO_SHAPE.NO_SYMBOL, Emu(mid - d // 2),
                                         Emu(lane - d // 2), Emu(d), Emu(d))
                sym.fill.background()
                sym.line.color.rgb = col
                sym.line.width = Pt(2.25)
                _flatten(sym)
            self._text(s, Emu(mid - int(Inches(1.6))), Emu(lane + int(Inches(0.16))),
                       Inches(3.2), Inches(0.4), "backlabel",
                       [{"t": backedge["label"], "size": TYPE["cap"], "color": col,
                         "bold": True, "align": PP_ALIGN.CENTER}])
        if legend:
            self._legend(s, legend)
        if source:
            self._footer(s, source)
        return s

    def hub(self, heading, core, spokes, kicker=None, caption=None, source=None):
        """One core node → N spoke nodes (same-core-different-hats)."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, heading)
        cx, cy, cw, ch = int(MARGIN), int(Inches(3.35)), int(Inches(2.9)), int(Inches(1.3))
        corebox = self._node(s, Emu(cx), Emu(cy), Emu(cw), Emu(ch), "action",
                             core, name="node")
        corebox.line.width = Pt(2.5)
        sx = int(Inches(6.4))
        sw = int(CONTENT_R) - sx
        m = len(spokes)
        sgap = int(Inches(0.35))
        sh = int((int(Inches(5.9)) - int(Inches(2.55)) - sgap * (m - 1)) / m)
        sy0 = int(Inches(2.55))
        for i, (sem, label, note) in enumerate(spokes):
            syi = sy0 + i * (sh + sgap)
            box = self._node(s, Emu(sx), Emu(syi), Emu(sw), Emu(sh), sem, label,
                             note, name="node")
            self._connect(s, corebox, box, 3, 1, color=TOKENS["grey"])
        if caption:
            self._text(s, Emu(cx), Emu(cy + ch + int(Inches(0.35))), Emu(cw + int(Inches(0.6))),
                       Inches(1.2), "caption",
                       [{"t": caption, "size": TYPE["body"], "color": TOKENS["grey"]}])
        if source:
            self._footer(s, source)
        return s

    def two_loops(self, heading, team, inner, kicker=None, source=None):
        """Hero nested-loop diagram: outer Team Loop containing an inner Agent
        Loop. team/inner are dicts (see build script for the shape)."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, heading)
        I = lambda v: int(Inches(v))

        # outer Team Loop box
        ox, oy, ow, oh = I(0.9), I(2.25), I(11.53), I(4.55)
        outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(ox), Emu(oy),
                                   Emu(ow), Emu(oh))
        _solid(outer, TOKENS["paper"], TOKENS["mist"], Pt(1.5))
        self._text(s, Emu(ox + I(0.28)), Emu(oy + I(0.16)), Emu(ow - I(0.5)),
                   Inches(0.4), "outer-label",
                   [{"t": team["outer"], "size": TYPE["cap"], "color": TOKENS["grey"],
                     "bold": True, "upper": True, "spc": 1.2}])

        # team row: steps + gate + done
        labels = team["steps"] + [("decision", team["gate"]), ("success", team["done"])]
        row_y, row_h = I(2.9), I(0.72)
        cells = self._row(len(labels), ox + I(0.3), ox + ow - I(0.3), I(0.22))
        tboxes = []
        for (x, w), (sem, lab) in zip(cells, labels):
            shape = MSO_SHAPE.DIAMOND if sem == "decision" else MSO_SHAPE.ROUNDED_RECTANGLE
            tboxes.append(self._node(s, Emu(x), Emu(row_y), Emu(w), Emu(row_h),
                                     sem, lab, name="node", shape=shape))
        for a, b in zip(tboxes, tboxes[1:]):
            self._connect(s, a, b, 3, 1, color=TOKENS["grey"])
        gate, done, plan = tboxes[-2], tboxes[-1], tboxes[0]
        # yes label on gate→done
        self._text(s, Emu(int(gate.left) + int(gate.width)), Emu(row_y - I(0.28)),
                   Inches(0.9), Inches(0.3), "yes",
                   [{"t": "yes", "size": TYPE["cap"], "color": SEM["success"][1],
                     "bold": True, "align": PP_ALIGN.CENTER}])

        # inner Agent Loop box (nested, lower-left)
        ix, iy, iw, ih = I(1.5), I(4.25), I(6.2), I(1.78)
        innerbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(ix), Emu(iy),
                                      Emu(iw), Emu(ih))
        _solid(innerbox, TOKENS["panel"], TOKENS["hairline"], Pt(1.25))
        self._text(s, Emu(ix + I(0.22)), Emu(iy + I(0.12)), Emu(iw - I(0.4)),
                   Inches(0.35), "inner-label",
                   [{"t": inner["label"], "size": TYPE["cap"], "color": TOKENS["grey"],
                     "bold": True, "upper": True, "spc": 1.2}])
        alabels = [("action", inner["steps"][0]), ("decision", inner["gate"]),
                   ("success", inner["done"])]
        arow_y, arow_h = I(4.95), I(0.62)
        acells = self._row(3, ix + I(0.28), ix + iw - I(0.28), I(0.22))
        aboxes = []
        for (x, w), (sem, lab) in zip(acells, alabels):
            shape = MSO_SHAPE.DIAMOND if sem == "decision" else MSO_SHAPE.ROUNDED_RECTANGLE
            aboxes.append(self._node(s, Emu(x), Emu(arow_y), Emu(w), Emu(arow_h),
                                     sem, lab, name="node", shape=shape))
        for a, b in zip(aboxes, aboxes[1:]):
            self._connect(s, a, b, 3, 1, color=TOKENS["grey"])
        # agent-loop return: gate → down → left → up to use-tools
        agate, atool = aboxes[1], aboxes[0]
        agx = int(agate.left) + int(agate.width) // 2
        atx = int(atool.left) + int(atool.width) // 2
        aby = arow_y + arow_h
        alane = iy + ih - I(0.2)
        self._polyline(s, [(agx, aby), (agx, alane), (atx, alane), (atx, aby)],
                       TOKENS["accent"], dash="dash")
        self._text(s, Emu(atx - I(0.2)), Emu(alane - I(0.02)), Inches(2.4), Inches(0.3),
                   "a-rerun", [{"t": inner["rerun"], "size": TYPE["cap"],
                                "color": TOKENS["accent"], "bold": True}])

        # each-member connector: team step → inner box top
        src = tboxes[team["from_idx"]]
        mx = int(src.left) + int(src.width) // 2
        self._polyline(s, [(mx, row_y + row_h), (mx, iy)], TOKENS["grey"])
        self._text(s, Emu(mx + I(0.14)), Emu(row_y + row_h + I(0.04)), Inches(3.0),
                   Inches(0.3), "each",
                   [{"t": team["each"], "size": TYPE["cap"], "color": TOKENS["grey"],
                     "italic": True}])

        # team-loop return: gate → down → far-left → up into plan (left side)
        gx = int(gate.left) + int(gate.width) // 2
        gby = row_y + row_h
        blane = oy + oh - I(0.28)
        plx = ox + I(0.15)
        ply = int(plan.top) + int(plan.height) // 2
        self._polyline(s, [(gx, gby), (gx, blane), (plx, blane), (plx, ply),
                           (int(plan.left), ply)], TOKENS["accent"], dash="dash")
        self._text(s, Emu(ix + I(0.1)), Emu(iy + ih + I(0.05)), Inches(4.6),
                   Inches(0.35), "t-rerun",
                   [{"t": team["rerun"], "size": TYPE["cap"], "color": TOKENS["accent"],
                     "bold": True}])
        if source:
            self._footer(s, source)
        return s

    # =====================================================================
    # MORE DIAGRAM ARCHETYPES (every slide is a picture)
    # =====================================================================
    _CHIP = {
        "accent": (TOKENS["accent"], TOKENS["paper"], None),
        "tint": (TOKENS["accentTint"], TOKENS["ink"], TOKENS["accent"]),
        "success": (SEM["success"][0], TOKENS["ink"], SEM["success"][1]),
        "panel": (TOKENS["panel"], TOKENS["ink"], TOKENS["hairline"]),
        "muted": (TOKENS["paper"], TOKENS["grey"], TOKENS["hairline"]),
    }

    def _chip(self, slide, x, y, w, h, label, tone="panel", align=PP_ALIGN.CENTER,
              size=None, name="chip"):
        fill, tcol, border = self._CHIP[tone]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(x)),
                                     Emu(int(y)), Emu(int(w)), Emu(int(h)))
        box.name = name
        _solid(box, fill, border, Pt(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.14)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        self._run(p, label, size or TYPE["cap"], tcol, bold=True)
        return box

    def _glyph(self, slide, x, y, size, kind, line_color, fill_color=None, w=Pt(3.0)):
        g = slide.shapes.add_shape(kind, Emu(int(x)), Emu(int(y)), Emu(int(size)),
                                   Emu(int(size)))
        if fill_color is not None:
            g.fill.solid()
            g.fill.fore_color.rgb = fill_color
        else:
            g.fill.background()
        g.line.color.rgb = line_color
        g.line.width = w
        _flatten(g)
        return g

    def transform(self, headline, before, after, kicker=None, source=None):
        """Before -> after: two panels with a big accent arrow between."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, headline)
        I = lambda v: int(Inches(v))
        py, ph, pw = I(2.85), I(3.15), I(5.1)
        lx = int(MARGIN)
        ax = lx + pw + I(0.1)
        aw = I(1.0)
        rx = ax + aw + I(0.1)

        def panel(x, spec, tone):
            fill, border, tcol = ((TOKENS["panel"], TOKENS["hairline"], TOKENS["grey"])
                                  if tone == "muted"
                                  else (TOKENS["accentTint"], TOKENS["accent"], TOKENS["ink"]))
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(py),
                                     Emu(pw), Emu(ph))
            _solid(box, fill, border, Pt(1.5))
            self._text(s, Emu(x + I(0.36)), Emu(py + I(0.3)), Emu(pw - I(0.7)),
                       Inches(0.4), "tag",
                       [{"t": spec["tag"], "size": TYPE["eye"],
                         "color": TOKENS["accent"] if tone == "tint" else TOKENS["grey"],
                         "bold": True, "upper": True, "spc": 1.6}])
            gs = I(1.25)
            self._glyph(s, x + pw // 2 - gs // 2, py + I(0.95), gs, spec["glyph"],
                        spec["gcolor"], spec.get("gfill"), w=Pt(3.5))
            self._text(s, Emu(x + I(0.4)), Emu(py + ph - I(0.95)), Emu(pw - I(0.8)),
                       Inches(0.8), "caption",
                       [{"t": spec["caption"], "size": TYPE["body"], "color": tcol,
                         "bold": True, "align": PP_ALIGN.CENTER}])

        panel(lx, before, "muted")
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(ax),
                                   Emu(py + ph // 2 - I(0.4)), Emu(aw), Emu(I(0.8)))
        _solid(arrow, TOKENS["accent"])
        panel(rx, after, "tint")
        if source:
            self._footer(s, source)
        return s

    def progression(self, headline, steps, kicker=None, source=None):
        """N connected nodes coloured by build status (built=green, new=accent)."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, headline)
        I = lambda v: int(Inches(v))
        n = len(steps)
        gap = I(0.5)
        nw = int((int(CONTENT_W) - gap * (n - 1)) / n)
        nh = I(1.9)
        y = I(3.5)
        x = int(MARGIN)
        boxes = []
        for st in steps:
            built = st["status"].lower().startswith("built")
            fill = SEM["success"][0] if built else TOKENS["accentTint"]
            border = SEM["success"][1] if built else TOKENS["accent"]
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y),
                                     Emu(nw), Emu(nh))
            box.name = "node"
            _solid(box, fill, border, Pt(2.0))
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.24)
            tf.margin_top = Inches(0.22)
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT
            self._run(p0, st["label"], TYPE["card"], TOKENS["ink"], bold=True)
            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.LEFT
            p1.space_before = Pt(7)
            self._run(p1, st["note"], TYPE["body"], TOKENS["ink"])
            _flatten(box)
            self._chip(s, x, y - I(0.62), I(1.5), I(0.44), st["status"],
                       tone="success" if built else "accent", name="badge")
            boxes.append(box)
            x += nw + gap
        for a, b in zip(boxes, boxes[1:]):
            self._connect(s, a, b, 3, 1, color=TOKENS["grey"])
        if source:
            self._footer(s, source)
        return s

    def gate_loop(self, headline, work, gate, done, backlabel, properties,
                  kicker=None, source=None):
        """A mini loop (work -> gate -> done, no -> back) + property callouts."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, headline)
        I = lambda v: int(Inches(v))
        y, h = I(2.55), I(1.0)
        specs = [("action", work), ("decision", gate), ("success", done)]
        cells = self._row(3, int(MARGIN) + I(1.3), int(CONTENT_R) - I(1.3), I(1.0))
        boxes = []
        for (x, w), (sem, lab) in zip(cells, specs):
            shape = MSO_SHAPE.DIAMOND if sem == "decision" else MSO_SHAPE.ROUNDED_RECTANGLE
            boxes.append(self._node(s, Emu(x), Emu(y), Emu(w), Emu(h), sem, lab,
                                    name="node", shape=shape))
        for a, b in zip(boxes, boxes[1:]):
            self._connect(s, a, b, 3, 1, color=TOKENS["grey"])
        self._text(s, Emu(int(boxes[1].left) + int(boxes[1].width)), Emu(y - I(0.3)),
                   Inches(0.8), Inches(0.3), "yes",
                   [{"t": "yes", "size": TYPE["cap"], "color": SEM["success"][1],
                     "bold": True, "align": PP_ALIGN.CENTER}])
        g, w0 = boxes[1], boxes[0]
        gx = int(g.left) + int(g.width) // 2
        wx = int(w0.left) + int(w0.width) // 2
        by = y + h
        lane = by + I(0.55)
        self._polyline(s, [(gx, by), (gx, lane), (wx, lane), (wx, by)],
                       TOKENS["accent"], dash="dash")
        self._text(s, Emu((gx + wx) // 2 - I(1.6)), Emu(lane + I(0.05)), Inches(3.2),
                   Inches(0.3), "backlabel",
                   [{"t": backlabel, "size": TYPE["cap"], "color": TOKENS["accent"],
                     "bold": True, "align": PP_ALIGN.CENTER}])
        # property callouts
        py = I(4.75)
        m = len(properties)
        gap = I(0.6)
        cw = int((int(CONTENT_W) - gap * (m - 1)) / m)
        x = int(MARGIN)
        for title, desc in properties:
            self._rule(s, Emu(x), Emu(py), Inches(0.7), TOKENS["accent"])
            self._text(s, Emu(x), Emu(py + I(0.26)), Emu(cw), Inches(0.6), "p-title",
                       [{"t": title, "size": TYPE["card"], "color": TOKENS["ink"],
                         "bold": True}])
            self._text(s, Emu(x), Emu(py + I(0.95)), Emu(cw), Inches(1.4), "p-desc",
                       [{"t": desc, "size": TYPE["body"], "color": TOKENS["grey"]}])
            x += cw + gap
        if source:
            self._footer(s, source)
        return s

    def timeline(self, headline, lanes, kicker=None, source=None):
        """Roadmap swimlanes: month header pills + stacked item chips, one
        lane marked as the milestone."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, headline)
        I = lambda v: int(Inches(v))
        n = len(lanes)
        gap = I(0.4)
        lw = int((int(CONTENT_W) - gap * (n - 1)) / n)
        x = int(MARGIN)
        top = I(2.45)
        base = top + I(0.95)
        # timeline baseline with arrow (time ->)
        self._polyline(s, [(int(MARGIN), base), (int(CONTENT_R), base)], TOKENS["mist"])
        for lane in lanes:
            milestone = lane["tone"] == "milestone"
            if milestone:
                bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x - I(0.12)),
                                        Emu(top - I(0.12)), Emu(lw + I(0.24)),
                                        Emu(I(4.35)))
                _solid(bg, TOKENS["accentTint"], TOKENS["accent"], Pt(1.25))
            # month pill
            self._chip(s, x, top, lw, I(0.55), lane["month"],
                       tone="accent" if milestone else "panel", size=TYPE["body"])
            if milestone:
                self._text(s, Emu(x), Emu(base + I(0.02)), Emu(lw), Inches(0.3),
                           "mstone", [{"t": "the milestone", "size": TYPE["cap"],
                                       "color": TOKENS["accent"], "bold": True,
                                       "italic": True, "align": PP_ALIGN.CENTER}])
            cy = base + I(0.45)
            for it in lane["items"]:
                self._chip(s, x + I(0.05), cy, lw - I(0.1), I(0.62), it,
                           tone="tint" if milestone else "muted",
                           align=PP_ALIGN.LEFT, size=TYPE["body"])
                cy += I(0.75)
            x += lw + gap
        if source:
            self._footer(s, source)
        return s

    def sequence(self, headline, left, right, kicker=None, source=None):
        """Two ordered horizons (August -> September) with a connecting arrow."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, headline)
        I = lambda v: int(Inches(v))
        py, ph, pw = I(2.7), I(3.3), I(5.1)
        lx = int(MARGIN)
        ax = lx + pw + I(0.1)
        aw = I(1.0)
        rx = ax + aw + I(0.1)

        def panel(x, data, tone):
            fill, border, tcol, bcol = ((TOKENS["panel"], TOKENS["hairline"],
                                         TOKENS["ink"], TOKENS["grey"]) if tone == "muted"
                                        else (TOKENS["accentTint"], TOKENS["accent"],
                                              TOKENS["ink"], TOKENS["ink"]))
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(py),
                                     Emu(pw), Emu(ph))
            _solid(box, fill, border, Pt(1.5))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.4)
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT
            self._run(p0, data["title"], TYPE["card"], tcol, bold=True)
            for it in data["items"]:
                pp = tf.add_paragraph()
                pp.alignment = PP_ALIGN.LEFT
                pp.space_before = Pt(10)
                self._run(pp, it, TYPE["body"], bcol)
            _flatten(box)

        panel(lx, left, "muted")
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(ax),
                                   Emu(py + ph // 2 - I(0.4)), Emu(aw), Emu(I(0.8)))
        _solid(arrow, TOKENS["accent"])
        panel(rx, right, "tint")
        if source:
            self._footer(s, source)
        return s

    def scope_boundary(self, headline, in_label, in_items, out_items, kicker=None,
                       source=None):
        """In-scope box vs deferred items marked with a no-symbol."""
        s = self._slide(bg=TOKENS["paper"])
        self._heading(s, kicker, headline)
        I = lambda v: int(Inches(v))
        jx, jy, jw, jh = int(MARGIN), I(2.7), I(4.9), I(3.6)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(jx), Emu(jy),
                                 Emu(jw), Emu(jh))
        _solid(box, TOKENS["accentTint"], TOKENS["accent"], Pt(2.0))
        self._text(s, Emu(jx + I(0.35)), Emu(jy + I(0.3)), Emu(jw - I(0.7)),
                   Inches(0.4), "in-label",
                   [{"t": in_label, "size": TYPE["eye"], "color": TOKENS["accent"],
                     "bold": True, "upper": True, "spc": 1.6}])
        cy = jy + I(1.0)
        for it in in_items:
            self._chip(s, jx + I(0.35), cy, jw - I(0.7), I(0.62), it, tone="accent",
                       align=PP_ALIGN.LEFT, size=TYPE["body"])
            cy += I(0.78)
        rx = jx + jw + I(0.7)
        rw = int(CONTENT_R) - rx
        self._text(s, Emu(rx), Emu(jy + I(0.02)), Emu(rw), Inches(0.4), "out-label",
                   [{"t": "not in July", "size": TYPE["eye"], "color": TOKENS["grey"],
                     "bold": True, "upper": True, "spc": 1.6}])
        oy = jy + I(0.75)
        for label, tag in out_items:
            gs = I(0.5)
            self._glyph(s, rx, oy, gs, MSO_SHAPE.NO_SYMBOL, SEM["error"][1], w=Pt(2.5))
            self._text(s, Emu(rx + I(0.7)), Emu(oy - I(0.02)), Emu(rw - I(0.7)),
                       Inches(0.6), "out-item",
                       [{"t": label, "size": TYPE["body"], "color": TOKENS["ink"],
                         "bold": True},
                        {"t": tag, "size": TYPE["cap"], "color": TOKENS["grey"]}])
            oy += I(0.92)
        if source:
            self._footer(s, source)
        return s

    def gantt(self, headline, months, bars, kicker=None, subtitle=None, source=None):
        """Roadmap Gantt: month columns shaded light->dark, horizontal bars that
        span their months. bars: [{label, start, end, phase}] where start/end are
        in month units [0..len(months)] and phase indexes the shade ramp."""
        s = self._slide(bg=TOKENS["paper"])
        I = lambda v: int(Inches(v))
        if kicker:
            self._kicker(s, kicker)
        self._text(s, MARGIN, Inches(1.15) if kicker else Inches(0.85), CONTENT_W,
                   Inches(0.7), "title",
                   [{"t": headline, "size": TYPE["h2"], "color": TOKENS["ink"],
                     "bold": True}])
        if subtitle:
            self._text(s, MARGIN, Inches(2.02), CONTENT_W, Inches(0.5), "subtitle",
                       [{"t": subtitle, "size": TYPE["sub"], "color": TOKENS["grey"]}])
        phases = GANTT_PHASES
        left, right = int(MARGIN), int(CONTENT_R)
        total = right - left
        n = len(months)
        col = total // n
        top, hh = I(2.55), I(0.55)
        for i, m in enumerate(months):
            x = left + i * col
            hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(top), Emu(col),
                                     Emu(hh))
            hdr.name = "hdr"
            _solid(hdr, phases[min(i, len(phases) - 1)])
            tf = hdr.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            self._run(p, m, TYPE["card"], TOKENS["paper"], bold=True)
            _flatten(hdr)
        gtop, gbot = top + hh, I(6.95)
        for i in range(1, n):
            x = left + i * col
            self._polyline(s, [(x, gtop), (x, gbot)], TOKENS["hairline"], arrow=False,
                           w=Pt(1.0))
        bar_top = gtop + I(0.3)
        slot = (gbot - bar_top) // max(len(bars), 1)
        bar_h = int(slot * 0.74)
        for j, bar in enumerate(bars):
            bx = left + int(bar["start"] / n * total)
            bw = max(int((bar["end"] - bar["start"]) / n * total), I(0.6))
            by = bar_top + j * slot
            r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(bx), Emu(by), Emu(bw),
                                   Emu(bar_h))
            r.name = "bar"
            _solid(r, phases[min(bar["phase"], len(phases) - 1)])
            tf = r.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.18)
            tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            self._run(p, bar["label"], TYPE["body"], TOKENS["paper"], bold=True)
            _flatten(r)
        if source:
            self._footer(s, source)
        return s
