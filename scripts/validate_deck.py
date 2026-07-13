#!/usr/bin/env python3
"""validate_deck.py — machine gate for generated .pptx decks (DK rules).

Run with a python that can import pptx (the vibe-pptx toolkit venv):
    ~/.claude/vibe-pptx/bin/python scripts/validate_deck.py <deck.pptx>

Prints one line per violation ("slide N: DK-x …") and exits 1 if any;
exits 0 with a single PASS line otherwise. Validates decks built by the
deck skill — explicit-everything is the contract (inherited/None font
sizes are violations, not unknowns).

Thresholds and the colour sets live here only, cited to DR numbers in the
deck-taste research report (.vibe/research/*deck-taste-anti-slop.md) and
mirrored in scripts/deck_lib.py + references/design-system.md — change them
in the research report first, then in lockstep.

Palette model (DR-5): chrome uses ONE accent over neutrals; diagrams may use
the curated semantic palette (SEM) and the roadmap Gantt ramp. So the
semantic/Gantt hues are *allowed* (DK-5) but do NOT count toward the ≤4-hue
"no rainbow chrome" limit (DK-7). The word budget (DK-1) counts prose only —
diagram/table labels (nodes, chips, table cells, …) are exempt, so a real
diagram or status matrix is never mistaken for a text wall.
"""
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

RULES = {
    "DK-1_max_words_per_slide": 50,      # DR-3 hard ceiling (prose only; target ~30)
    "DK-2_min_font_pt": 14,              # DR-4 absolute floor (body is 18pt by token)
    "DK-3_bounds_tolerance_emu": 9525,   # DR-7 on-slide bounds, 1px @96dpi tolerance
    "DK-4_overlap_tolerance_emu2": 91440 * 91440,  # DR-7 ~1% sq-inch overlap
    "DK-6_max_title_words": 15,          # DR-2 assertion title ceiling
    "DK-7_max_hues_per_slide": 4,        # DR-5 non-neutral, non-semantic hues
}

# --- colour sets (DR-5), mirroring deck_lib TOKENS / SEM / GANTT_PHASES ------
NEUTRALS = {"FFFFFF", "16283F", "EEF2F7", "5C6675", "93A1B2", "DCE3EC",
            "F5F5F5", "666666"}                 # structure: bg, text, panels, rules
ACCENTS = {"1F5FA8", "D7E3F1"}                  # the ONE accent family (counted by DK-7)
SEM = {"DAE8FC", "6C8EBF", "D5E8D4", "82B366",  # semantic diagram swatches
       "F8CECC", "B85450", "FFF2CC", "D6B656",
       "E1D5E7", "9673A6"}
GANTT = {"3E7CBF", "1F5FA8", "123A66"}          # roadmap month ramp
ALLOWED = NEUTRALS | ACCENTS | SEM | GANTT      # DK-5: every hex must be in here
# hues exempt from the DK-7 rainbow count (semantic + non-accent Gantt shades)
DK7_EXEMPT = SEM | (GANTT - ACCENTS)

# DK-1 exemption: shapes whose text is a diagram/table label, not prose.
DIAGRAM_PREFIXES = (
    "node", "chip", "badge", "cell", "legend", "swatch", "bar", "hdr",
    "outer", "inner", "each", "yes", "backlabel", "rerun", "a-rerun",
    "t-rerun", "mstone", "in-label", "out-label", "out-item", "tag",
    "pt-", "p-title", "p-desc", "kicker", "mark",
)
DASHES = ("—", "–")  # DK-8: em / en dash (the #1 AI tell)


def _rgb(color):
    # python-pptx raises AttributeError/TypeError for non-RGB color types
    # (theme colors, inherited); those are "no explicit RGB" here, by design.
    try:
        if color and color.type is not None and color.rgb is not None:
            return str(color.rgb)
    except (AttributeError, TypeError):
        return None
    return None


def _is_connector(shape):
    # python-pptx reports a connector's shape_type as MSO_SHAPE_TYPE.LINE.
    return shape.shape_type == MSO_SHAPE_TYPE.LINE


def _rect(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _overlap_area(a, b):
    w = min(a[2], b[2]) - max(a[0], b[0])
    h = min(a[3], b[3]) - max(a[1], b[1])
    return w * h if (w > 0 and h > 0) else 0


def validate(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    tol = RULES["DK-3_bounds_tolerance_emu"]
    violations = []

    for idx, slide in enumerate(prs.slides, 1):
        words = 0
        hues = set()
        text_rects = []
        for shape in slide.shapes:
            if shape.left is None:
                continue
            l, t, r, b = _rect(shape)
            if l < -tol or t < -tol or r > sw + tol or b > sh + tol:
                violations.append(
                    f"slide {idx}: DK-3 shape '{shape.name}' extends off-slide")
            if _is_connector(shape):
                continue
            fill_rgb = None
            try:
                if shape.fill.type == 1:  # MSO_FILL.SOLID
                    fill_rgb = _rgb(shape.fill.fore_color)
            except (AttributeError, TypeError, NotImplementedError):
                fill_rgb = None
            if fill_rgb and fill_rgb not in NEUTRALS:
                if fill_rgb not in ALLOWED:
                    violations.append(
                        f"slide {idx}: DK-5 fill #{fill_rgb} on '{shape.name}' not in palette")
                if fill_rgb not in DK7_EXEMPT:
                    hues.add(fill_rgb)
            if not shape.has_text_frame:
                continue
            is_diagram = shape.name.lower().startswith(DIAGRAM_PREFIXES)
            has_text = False
            shape_words = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text.strip()
                    if not txt:
                        continue
                    has_text = True
                    shape_words += len(txt.split())
                    if any(d in run.text for d in DASHES):
                        violations.append(
                            f"slide {idx}: DK-8 em/en-dash in '{shape.name}' (use a plain hyphen)")
                    if run.font.size is None:
                        violations.append(
                            f"slide {idx}: DK-2 run without explicit font size in '{shape.name}'")
                    elif run.font.size.pt < RULES["DK-2_min_font_pt"]:
                        violations.append(
                            f"slide {idx}: DK-2 font {run.font.size.pt:.0f}pt < "
                            f"{RULES['DK-2_min_font_pt']}pt in '{shape.name}'")
                    run_rgb = _rgb(run.font.color)
                    if run_rgb is None:
                        violations.append(
                            f"slide {idx}: DK-2 run without explicit font color in '{shape.name}'")
                    elif run_rgb not in NEUTRALS:
                        if run_rgb not in ALLOWED:
                            violations.append(
                                f"slide {idx}: DK-5 font #{run_rgb} on '{shape.name}' not in palette")
                        if run_rgb not in DK7_EXEMPT:
                            hues.add(run_rgb)
            if has_text and not is_diagram:
                words += shape_words
            if has_text:
                text_rects.append((shape.name, _rect(shape), is_diagram))
            if shape.name.lower().startswith("title"):
                title_words = len(shape.text_frame.text.split())
                if title_words > RULES["DK-6_max_title_words"]:
                    violations.append(
                        f"slide {idx}: DK-6 title has {title_words} words > "
                        f"{RULES['DK-6_max_title_words']}")
        if words > RULES["DK-1_max_words_per_slide"]:
            violations.append(
                f"slide {idx}: DK-1 {words} prose words > {RULES['DK-1_max_words_per_slide']}")
        if len(hues) > RULES["DK-7_max_hues_per_slide"]:
            violations.append(
                f"slide {idx}: DK-7 {len(hues)} non-neutral chrome hues > "
                f"{RULES['DK-7_max_hues_per_slide']}")
        for i in range(len(text_rects)):
            for j in range(i + 1, len(text_rects)):
                if text_rects[i][2] or text_rects[j][2]:
                    continue  # diagram labels: placed by layout, render-reviewed
                if _overlap_area(text_rects[i][1], text_rects[j][1]) > RULES["DK-4_overlap_tolerance_emu2"]:
                    violations.append(
                        f"slide {idx}: DK-4 text shapes '{text_rects[i][0]}' and "
                        f"'{text_rects[j][0]}' overlap")
    return violations


def main():
    if len(sys.argv) != 2:
        print("usage: validate_deck.py <deck.pptx>", file=sys.stderr)
        return 2
    violations = validate(sys.argv[1])
    for v in violations:
        print(v)
    if violations:
        return 1
    print("DECK PASS — all DK rules green")
    return 0


if __name__ == "__main__":
    sys.exit(main())
