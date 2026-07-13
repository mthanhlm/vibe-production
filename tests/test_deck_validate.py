"""Tests for scripts/validate_deck.py — run through a pptx-capable python.

Resolution order for the interpreter: the vibe-pptx toolkit venv if
present, else the current python when it can import pptx (CI installs
python-pptx), else the whole module is skipped cleanly. The validator is
always exercised as a subprocess, never imported.
"""
import os
import subprocess
import sys
import tempfile
import textwrap
import unittest

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
VALIDATOR = os.path.join(REPO, "scripts", "validate_deck.py")
VENV_PY = os.path.expanduser(
    os.path.join(os.environ.get("VIBE_PPTX_DIR", "~/.claude/vibe-pptx"),
                 "bin", "python"))


def _pptx_python():
    if os.path.exists(VENV_PY):
        return VENV_PY
    probe = subprocess.run([sys.executable, "-c", "import pptx"],
                           capture_output=True)
    return sys.executable if probe.returncode == 0 else None


PPTX_PY = _pptx_python()

# Colours are Ink & Signal tokens (see deck_lib.TOKENS): ink 16283F, grey 5C6675.
GOOD_DECK = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(11), Inches(1.2))
tb.name = "title 1"
r = tb.text_frame.paragraphs[0].add_run(); r.text = "Revenue doubled in one quarter"
r.font.size = Pt(36); r.font.color.rgb = RGBColor(0x16, 0x28, 0x3F)
b = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11), Inches(2))
r = b.text_frame.paragraphs[0].add_run(); r.text = "Six words of supporting body text"
r.font.size = Pt(18); r.font.color.rgb = RGBColor(0x5C, 0x66, 0x75)
prs.save(OUT)
"""

BAD_DECK = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
# DK-2: tiny font; DK-5: off-palette color; DK-1: word wall (unnamed box counts)
tb = s.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(11), Inches(4))
r = tb.text_frame.paragraphs[0].add_run()
r.text = " ".join(["filler"] * 200)
r.font.size = Pt(9); r.font.color.rgb = RGBColor(0xFF, 0x00, 0xFF)
# DK-3: off-slide shape
off = s.shapes.add_textbox(Inches(15), Inches(1), Inches(3), Inches(1))
r2 = off.text_frame.paragraphs[0].add_run(); r2.text = "off slide"
r2.font.size = Pt(18)
# DK-4: two overlapping text boxes
o1 = s.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
r3 = o1.text_frame.paragraphs[0].add_run(); r3.text = "overlap one"
r3.font.size = Pt(18)
o2 = s.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(4), Inches(1))
r4 = o2.text_frame.paragraphs[0].add_run(); r4.text = "overlap two"
r4.font.size = Pt(18)
prs.save(OUT)
"""


# Exercises EVERY deck_lib layout — text-first and the semantic-diagram set
# (flow with a crossed back-edge, two_loops, gate_loop, progression, matrix,
# transform, sequence, scope_boundary, timeline, gantt, hub). Regresses:
#  - the multi-hue semantic palette must pass DK-7 (exempted), and the matrix /
#    diagram-label word volume must pass DK-1 (prose-only), i.e. a real diagram
#    is never mistaken for a text wall / rainbow;
#  - every shape stays shadowless (no <p:style> survives).
FULL_DECK = """
import sys
sys.path.insert(0, PLUGIN_SCRIPTS)
from deck_lib import Deck, TOKENS
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
d = Deck()
d.title("Cover assertion here", subtitle="sub", eyebrow="Internal")
d.section("01", "Section label")
d.statement("We can open three locations in 2026 from cash flow.")
d.title_body("Three forces point the same way",
             ["Two of four sites hit capacity daily.",
              "Catering has a standing waitlist."], eyebrow="Demand")
d.big_number("92%", "of revenue is repeat customers",
             title="Loyalty carries the business", source="Source: internal")
d.cards("What changes", [("A", "one"), ("B", "two"), ("C", "three")])
d.flow("How a change reaches production",
       ["Commit", "Build", "Verify", "Promote"], eyebrow="Pipeline")
d.flow("Today it never loops back",
       [("external", "Ask"), ("action", "Plan"), ("action", "Ship")],
       backedge={"label": "no feedback", "crossed": True}, kicker="Problem")
d.hub("One core, many roles", core="Core",
      spokes=[("action", "Plan", "a"), ("action", "Merge", "b"),
              ("action", "Check", "c")], kicker="Roles", caption="One core.")
d.two_loops("Two loops until the work is done",
            team={"outer": "Team loop", "steps": [("action", "Plan"), ("action", "Run")],
                  "gate": "Met?", "done": "Done", "rerun": "no", "from_idx": 1, "each": "each"},
            inner={"label": "Agent loop", "steps": ["Use tools"], "gate": "done?",
                   "done": "Return", "rerun": "no"}, kicker="Milestone")
d.gate_loop("The check is the only exit", work="Do work", gate="Met?", done="Ship",
            backlabel="no, retry",
            properties=[("Set here", "by the lead."), ("First seen", "each turn."),
                        ("Only exit", "when met.")], kicker="Contract")
d.progression("Three deliverables build up",
              [{"label": "A", "note": "one.", "status": "Built"},
               {"label": "B", "note": "two.", "status": "New"},
               {"label": "C", "note": "three.", "status": "New"}], kicker="Ships")
d.matrix("Status by item", ["Item", "Today", "End"],
         [["Alpha", "built", "hardened"], ["Beta", "no", "new"], ["Gamma", "no", "new"]],
         emphasize={0}, kicker="Status", col_weights=[2, 1, 1])
d.transform("From one pass to a loop",
            before={"tag": "Today", "glyph": MSO_SHAPE.NO_SYMBOL, "gcolor": TOKENS["grey"],
                    "caption": "stops."},
            after={"tag": "Later", "glyph": MSO_SHAPE.CIRCULAR_ARROW, "gcolor": TOKENS["accent"],
                   "caption": "loops."}, kicker="Shift")
d.sequence("Two horizons ahead",
           left={"title": "Aug", "items": ["one", "two"]},
           right={"title": "Sep", "items": ["three", "four"]}, kicker="Next")
d.scope_boundary("Narrow and deep on purpose", in_label="In scope",
                 in_items=["Alpha", "Beta"],
                 out_items=[("Gamma", "later"), ("Delta", "later")], kicker="Guardrail")
d.timeline("Roadmap by month",
           [{"month": "JUL", "tone": "milestone", "items": ["Alpha", "Beta"]},
            {"month": "AUG", "tone": "muted", "items": ["Gamma"]},
            {"month": "SEP", "tone": "muted", "items": ["Delta"]}], kicker="Plan")
d.gantt("Roadmap across the quarter", ["JUL", "AUG", "SEP"],
        [{"label": "Alpha", "start": 0, "end": 1, "phase": 0},
         {"label": "Beta", "start": 0.5, "end": 2, "phase": 1},
         {"label": "Gamma", "start": 1.5, "end": 3, "phase": 2}],
        kicker="Plan", subtitle="Q3 2026")
d.closing("Questions?", contact="ops@example.com")
d.save(OUT)
# shadow regression: every deck_lib shape is flattened, so no slide keeps a style ref
styles = 0
for slide in d.prs.slides:
    styles += len(slide._element.findall('.//' + qn('p:style')))
if styles:
    sys.exit("REGRESSION: %d <p:style> survived (drop-shadow)" % styles)
"""


@unittest.skipUnless(PPTX_PY, "no python-pptx interpreter available")
class TestDeckLib(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="vibe-decklib-test.")
        self.addCleanup(lambda: subprocess.run(["rm", "-rf", self.tmp]))

    def test_full_deck_all_layouts_valid_and_shadowless(self):
        out = os.path.join(self.tmp, "full.pptx")
        script = (f"OUT = {out!r}\nPLUGIN_SCRIPTS = "
                  f"{os.path.join(REPO, 'scripts')!r}\n" + textwrap.dedent(FULL_DECK))
        built = subprocess.run([PPTX_PY, "-c", script],
                               capture_output=True, text=True, timeout=90)
        self.assertEqual(built.returncode, 0, built.stderr)  # incl. shadow regression
        validated = subprocess.run([PPTX_PY, VALIDATOR, out],
                                   capture_output=True, text=True, timeout=90)
        self.assertEqual(validated.returncode, 0,
                         "a shipped layout regressed:\n" + validated.stdout)
        self.assertIn("DECK PASS", validated.stdout)


@unittest.skipUnless(PPTX_PY, "no python-pptx interpreter available")
class TestValidateDeck(unittest.TestCase):
    def _build(self, body, name):
        path = os.path.join(self.tmp, name)
        script = f"OUT = {path!r}\n" + textwrap.dedent(body)
        proc = subprocess.run([PPTX_PY, "-c", script],
                              capture_output=True, text=True, timeout=60)
        self.assertEqual(proc.returncode, 0, proc.stderr)
        return path

    def _validate(self, deck):
        return subprocess.run([PPTX_PY, VALIDATOR, deck],
                              capture_output=True, text=True, timeout=60)

    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="vibe-deck-test.")
        self.addCleanup(lambda: subprocess.run(["rm", "-rf", self.tmp]))

    def test_clean_deck_passes(self):
        proc = self._validate(self._build(GOOD_DECK, "good.pptx"))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)
        self.assertIn("DECK PASS", proc.stdout)

    def test_broken_deck_fails_with_rule_ids(self):
        proc = self._validate(self._build(BAD_DECK, "bad.pptx"))
        self.assertEqual(proc.returncode, 1)
        for rule in ("DK-1", "DK-2", "DK-3", "DK-4", "DK-5"):
            self.assertIn(rule, proc.stdout, f"{rule} missing:\n{proc.stdout}")
        self.assertTrue(all("slide 1:" in l for l in
                            proc.stdout.splitlines() if l.startswith("slide")))

    def test_usage_error(self):
        proc = subprocess.run([PPTX_PY, VALIDATOR],
                              capture_output=True, text=True, timeout=60)
        self.assertEqual(proc.returncode, 2)

    def test_colored_connector_not_palette_flagged(self):
        # Regresses the _is_connector exemption: an accent-line connector
        # must be skipped by the palette/bounds checks, not flagged.
        deck = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
a = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1))
a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xEE,0xF2,0xF7); a.line.color.rgb = RGBColor(0x1F,0x5F,0xA8)
b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(3), Inches(2), Inches(1))
b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xEE,0xF2,0xF7); b.line.color.rgb = RGBColor(0x1F,0x5F,0xA8)
c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
c.begin_connect(a, 3); c.end_connect(b, 1); c.line.color.rgb = RGBColor(0x1F,0x5F,0xA8)
prs.save(OUT)
"""
        proc = self._validate(self._build(deck, "conn.pptx"))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)

    def test_semantic_diagram_palette_passes(self):
        # The curated semantic node palette (blue/yellow/green/red) is many
        # hues but must PASS: DK-5 allows it, DK-7 exempts it from the rainbow
        # count. Regresses "diagrams look like slop to the validator".
        deck = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
for i, (fill, x) in enumerate([("DAE8FC",1),("FFF2CC",4),("D5E8D4",7),("F8CECC",10)]):
    n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3), Inches(2), Inches(1))
    n.name = "node %d" % i
    n.fill.solid(); n.fill.fore_color.rgb = RGBColor.from_string(fill)
    r = n.text_frame.paragraphs[0].add_run(); r.text = "N"
    r.font.size = Pt(16); r.font.color.rgb = RGBColor(0x16,0x28,0x3F)
prs.save(OUT)
"""
        proc = self._validate(self._build(deck, "sem.pptx"))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)

    def test_em_dash_flagged_dk8(self):
        deck = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(9), Inches(1))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "days — hours"
r.font.size = Pt(18); r.font.color.rgb = RGBColor(0x16,0x28,0x3F)
prs.save(OUT)
"""
        proc = self._validate(self._build(deck, "dash.pptx"))
        self.assertEqual(proc.returncode, 1)
        self.assertIn("DK-8", proc.stdout)

    def test_tiny_overlap_within_tolerance_passes(self):
        # ~0.5% sq-in sliver overlap must pass DK-4 (regresses the old
        # 100x-too-strict tolerance that hard-failed sub-pixel touches).
        deck = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
a = s.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
r = a.text_frame.paragraphs[0].add_run(); r.text = "box a"; r.font.size = Pt(18); r.font.color.rgb = RGBColor(0x16,0x28,0x3F)
b = s.shapes.add_textbox(Inches(3.95), Inches(3.9), Inches(3), Inches(0.5))
r = b.text_frame.paragraphs[0].add_run(); r.text = "box b"; r.font.size = Pt(18); r.font.color.rgb = RGBColor(0x16,0x28,0x3F)
prs.save(OUT)
"""
        proc = self._validate(self._build(deck, "sliver.pptx"))
        self.assertEqual(proc.returncode, 0, proc.stdout + proc.stderr)

    def test_missing_font_color_flagged_dk2(self):
        deck = """
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "no explicit color"
r.font.size = Pt(20)  # size set, color left inherited
prs.save(OUT)
"""
        proc = self._validate(self._build(deck, "nocolor.pptx"))
        self.assertEqual(proc.returncode, 1)
        self.assertIn("DK-2", proc.stdout)
        self.assertIn("font color", proc.stdout)


if __name__ == "__main__":
    unittest.main()
